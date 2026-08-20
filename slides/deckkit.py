#!/usr/bin/env python3
"""
deckkit.py - a small professional slide-design system on top of python-pptx.

It gives every bootcamp module the same clean look: a colored header band, a
consistent footer, and a set of *graphic* slide builders (timelines, chevron
flows, pipelines, card grids, comparison panels, KPI tiles, bar charts, and
takeaways) so the decks are visual, not walls of text.

House rules enforced automatically by clean_text():
  - no em dashes / en dashes (converted to a plain hyphen)
  - no emojis or pictographs (stripped)

Usage:
    from deckkit import Deck, PALETTE
    d = Deck(module_no=1, title="Foundations", subtitle="AI for Blue Team Operations",
             accent=PALETTE["blue"])
    d.title_slide(day="Day 1", duration="2 hours")
    d.timeline_slide(...)
    ...
    d.save("Module1_Foundations.pptx")
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Palette (brand-consistent; each module picks an accent for variety)
# ---------------------------------------------------------------------------
PALETTE = {
    "purple":  RGBColor(0x5B, 0x2A, 0xA5),   # Valix brand
    "purple_d":RGBColor(0x3B, 0x1A, 0x6E),
    "gold":    RGBColor(0xF2, 0xC0, 0x1E),    # Evolve accent
    "ink":     RGBColor(0x15, 0x16, 0x22),
    "slate":   RGBColor(0x3A, 0x3F, 0x4B),
    "muted":   RGBColor(0x6B, 0x72, 0x80),
    "cloud":   RGBColor(0xF4, 0xF2, 0xFB),    # light panel
    "cloud2":  RGBColor(0xEC, 0xEF, 0xF5),
    "white":   RGBColor(0xFF, 0xFF, 0xFF),
    "line":    RGBColor(0xD7, 0xDA, 0xE2),
    "blue":    RGBColor(0x25, 0x63, 0xEB),    # M1
    "teal":    RGBColor(0x0E, 0x9F, 0x8E),    # M2
    "violet":  RGBColor(0x6D, 0x28, 0xD9),    # M3
    "red":     RGBColor(0xDC, 0x26, 0x2A),    # M4
    "amber":   RGBColor(0xB4, 0x53, 0x09),    # M5
    "green":   RGBColor(0x1F, 0x9D, 0x6B),
    "danger":  RGBColor(0xE5, 0x48, 0x4D),
    "danger_bg":RGBColor(0xFC, 0xEC, 0xEC),
    "green_bg":RGBColor(0xE9, 0xF7, 0xF0),
    "blue_bg": RGBColor(0xEA, 0xF1, 0xFE),
}

TOPCAP = MSO_SHAPE.ROUND_2_SAME_RECTANGLE  # rounds only the top two corners

FONT = "Calibri"
FONT_H = "Calibri"

SW = Inches(13.333)
SH = Inches(7.5)


def clean_text(s) -> str:
    """Enforce house rules: no em/en dashes, no emoji/pictographs."""
    if s is None:
        return ""
    s = str(s)
    for bad in ("—", "–", "―", "‒"):
        s = s.replace(bad, "-")
    s = s.replace("→", "to").replace("’", "'").replace("‘", "'")
    s = s.replace("“", '"').replace("”", '"').replace("…", "...")
    out = []
    for ch in s:
        o = ord(ch)
        if (0x1F000 <= o <= 0x1FAFF) or (0x2600 <= o <= 0x27BF) or \
           (0xFE00 <= o <= 0xFE0F) or (0x1F1E6 <= o <= 0x1F1FF) or \
           o in (0x2705, 0x274C, 0x2714, 0x2716, 0x2B50, 0x2757, 0x2764, 0x2049, 0x203C):
            continue
        out.append(ch)
    return "".join(out)


class Deck:
    def __init__(self, module_no: int, title: str, subtitle: str, accent: RGBColor):
        self.module_no = module_no
        self.title = title
        self.subtitle = subtitle
        self.accent = accent
        self.prs = Presentation()
        self.prs.slide_width = SW
        self.prs.slide_height = SH
        self._blank = self.prs.slide_layouts[6]

    # ---- low-level helpers -------------------------------------------------
    def _slide(self):
        return self.prs.slides.add_slide(self._blank)

    def _rect(self, slide, x, y, w, h, fill=None, line=None, line_w=1.0,
              shape=MSO_SHAPE.RECTANGLE, shadow=False):
        sp = slide.shapes.add_shape(shape, x, y, w, h)
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        if shadow:
            self._soft_shadow(sp)
        return sp

    def _soft_shadow(self, sp):
        # add a subtle outer shadow via raw XML
        spPr = sp._element.spPr
        effLst = spPr.makeelement(qn('a:effectLst'), {})
        outer = spPr.makeelement(qn('a:outerShdw'),
                                 {'blurRad': '60000', 'dist': '25000', 'dir': '5400000',
                                  'rotWithShape': '0'})
        clr = spPr.makeelement(qn('a:srgbClr'), {'val': '1A1A2E'})
        alpha = spPr.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(alpha)
        outer.append(clr)
        effLst.append(outer)
        spPr.append(effLst)

    def _text(self, shape, blocks, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              wrap=True):
        """blocks = list of (text, size, bold, color, [space_after_pt]) tuples."""
        tf = shape.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)
        for i, blk in enumerate(blocks):
            text, size, bold, color = blk[0], blk[1], blk[2], blk[3]
            space_after = blk[4] if len(blk) > 4 else 4
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(space_after)
            p.space_before = Pt(0)
            r = p.add_run()
            r.text = clean_text(text)
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
        return tf

    def _textbox(self, slide, x, y, w, h, blocks, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(x, y, w, h)
        self._text(tb, blocks, align=align, anchor=anchor)
        return tb

    def _chrome(self, slide, header, kicker=None, band=True):
        """Add the standard header band + footer to a content slide."""
        if band:
            self._rect(slide, 0, 0, SW, Inches(1.15), fill=self.accent)
            self._rect(slide, 0, Inches(1.15), SW, Inches(0.06), fill=PALETTE["gold"])
            if kicker:
                self._textbox(slide, Inches(0.55), Inches(0.14), Inches(11), Inches(0.3),
                              [(kicker.upper(), 11, True, PALETTE["gold"])])
            self._textbox(slide, Inches(0.55), Inches(0.40), Inches(12.2), Inches(0.66),
                          [(header, 26, True, PALETTE["white"])], anchor=MSO_ANCHOR.MIDDLE)
        # footer
        self._rect(slide, 0, Inches(7.12), SW, Inches(0.38), fill=PALETTE["cloud2"])
        self._textbox(slide, Inches(0.4), Inches(7.14), Inches(8), Inches(0.32),
                      [("Valix AI  x  Evolve Academy   |   AI Blue Team Bootcamp", 9, False,
                        PALETTE["muted"])], anchor=MSO_ANCHOR.MIDDLE)
        n = len(self.prs.slides._sldIdLst)
        self._textbox(slide, Inches(11.4), Inches(7.14), Inches(1.5), Inches(0.32),
                      [("Module %d  |  %02d" % (self.module_no, n), 9, False, PALETTE["muted"])],
                      align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    def _chip(self, slide, x, y, text, fill, txtcolor=None, w=None):
        w = w or Inches(1.6)
        c = self._rect(slide, x, y, w, Inches(0.34), fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        c.adjustments[0] = 0.5
        self._text(c, [(text, 10.5, True, txtcolor or PALETTE["white"])],
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return c

    # ---- slide builders ----------------------------------------------------
    def title_slide(self, day: str, duration: str, tagline: str = ""):
        slide = self._slide()
        self._rect(slide, 0, 0, SW, SH, fill=PALETTE["ink"])
        # accent geometry
        self._rect(slide, 0, 0, Inches(0.28), SH, fill=self.accent)
        self._rect(slide, Inches(0.28), 0, Inches(0.08), SH, fill=PALETTE["gold"])
        band = self._rect(slide, Inches(9.7), 0, Inches(3.633), SH, fill=self.accent)
        self._soft_shadow(band)
        self._rect(slide, Inches(9.62), 0, Inches(0.08), SH, fill=PALETTE["gold"])
        # module number big on the accent band
        self._textbox(slide, Inches(9.7), Inches(2.5), Inches(3.633), Inches(2),
                      [("MODULE", 16, True, PALETTE["white"]),
                       (str(self.module_no), 120, True, PALETTE["white"])],
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # kicker
        self._textbox(slide, Inches(0.85), Inches(1.5), Inches(8.4), Inches(0.5),
                      [("AI BLUE TEAM AND INTRO TO AI RED TEAMING BOOTCAMP", 13, True,
                        PALETTE["gold"])])
        # title
        self._textbox(slide, Inches(0.8), Inches(2.4), Inches(8.6), Inches(1.7),
                      [(self.title, 46, True, PALETTE["white"])])
        self._textbox(slide, Inches(0.85), Inches(3.9), Inches(8.4), Inches(0.9),
                      [(self.subtitle, 22, False, PALETTE["line"])])
        if tagline:
            self._textbox(slide, Inches(0.85), Inches(4.75), Inches(8.4), Inches(0.7),
                          [(tagline, 14, False, PALETTE["muted"])])
        # info chips
        self._chip(slide, Inches(0.85), Inches(5.7), day, self.accent, w=Inches(1.5))
        self._chip(slide, Inches(2.5), Inches(5.7), duration, PALETTE["slate"], w=Inches(1.6))
        self._chip(slide, Inches(4.25), Inches(5.7), "Hands-on labs", PALETTE["slate"], w=Inches(2.0))
        self._textbox(slide, Inches(0.85), Inches(6.55), Inches(8), Inches(0.4),
                      [("Presenter: Dr. Gaspard Baye, Valix AI", 12, False, PALETTE["muted"])])
        return slide

    def section_slide(self, kicker, big_title, sub=""):
        slide = self._slide()
        self._rect(slide, 0, 0, SW, SH, fill=self.accent)
        self._rect(slide, 0, Inches(3.35), SW, Inches(0.06), fill=PALETTE["gold"])
        self._textbox(slide, Inches(1.0), Inches(2.4), Inches(11.3), Inches(0.5),
                      [(kicker.upper(), 14, True, PALETTE["gold"])], align=PP_ALIGN.CENTER)
        self._textbox(slide, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.2),
                      [(big_title, 40, True, PALETTE["white"])], align=PP_ALIGN.CENTER)
        if sub:
            self._textbox(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8),
                          [(sub, 16, False, PALETTE["cloud"])], align=PP_ALIGN.CENTER)
        return slide

    def bullets_slide(self, header, points, kicker=None, intro=None):
        slide = self._slide()
        self._chrome(slide, header, kicker)
        y = 1.5
        if intro:
            self._textbox(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.6),
                          [(intro, 15, False, PALETTE["slate"])])
            y = 2.1
        top = Inches(y)
        for i, pt in enumerate(points):
            head, body = (pt if isinstance(pt, tuple) else (pt, None))
            row = self._rect(slide, Inches(0.6), top, Inches(12.1), Inches(0.72),
                             fill=PALETTE["cloud"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            row.adjustments[0] = 0.12
            dot = self._rect(slide, Inches(0.78), top + Inches(0.16), Inches(0.4), Inches(0.4),
                             fill=self.accent, shape=MSO_SHAPE.OVAL)
            self._text(dot, [(str(i + 1), 14, True, PALETTE["white"])],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            blocks = [(head, 15, True, PALETTE["ink"])]
            if body:
                blocks.append((body, 12, False, PALETTE["slate"]))
            self._text(self._add_tb(slide, Inches(1.4), top + Inches(0.04),
                                    Inches(11.1), Inches(0.64)), blocks,
                       anchor=MSO_ANCHOR.MIDDLE)
            top += Inches(0.84)
        return slide

    def _add_tb(self, slide, x, y, w, h):
        return slide.shapes.add_textbox(x, y, w, h)

    def cards_slide(self, header, cards, columns=3, kicker=None, intro=None):
        """cards = list of (title, body) or (title, body, color)."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        top = 1.55
        if intro:
            self._textbox(slide, Inches(0.6), Inches(1.35), Inches(12.1), Inches(0.5),
                          [(intro, 14, False, PALETTE["slate"])])
            top = 2.05
        n = len(cards)
        rows = (n + columns - 1) // columns
        gap = 0.3
        margin = 0.6
        cw = (13.333 - 2 * margin - (columns - 1) * gap) / columns
        avail_h = 7.0 - top
        ch = (avail_h - (rows - 1) * gap) / rows
        ch = min(ch, 2.4)
        for idx, card in enumerate(cards):
            r = idx // columns
            c = idx % columns
            title = card[0]
            body = card[1]
            color = card[2] if len(card) > 2 else self.accent
            x = Inches(margin + c * (cw + gap))
            y = Inches(top + r * (ch + gap))
            box = self._rect(slide, x, y, Inches(cw), Inches(ch), fill=PALETTE["white"],
                             line=PALETTE["line"], shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                             shadow=True)
            box.adjustments[0] = 0.06
            self._rect(slide, x, y, Inches(cw), Inches(0.14), fill=color,
                       shape=TOPCAP)
            tb = self._add_tb(slide, x + Inches(0.2), y + Inches(0.28),
                              Inches(cw - 0.4), Inches(ch - 0.42))
            self._text(tb, [(title, 15, True, color),
                            (body, 12, False, PALETTE["slate"])])
        return slide

    def flow_slide(self, header, steps, kicker=None, caption=None, colors=None):
        """Chevron process flow. steps = list of (label, sublabel)."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        n = len(steps)
        margin = 0.5
        total_w = 13.333 - 2 * margin
        overlap = 0.28
        cw = (total_w + (n - 1) * overlap) / n
        y = Inches(2.7)
        h = Inches(1.7)
        for i, step in enumerate(steps):
            label, sub = (step if isinstance(step, tuple) else (step, None))
            color = (colors[i] if colors else self.accent)
            x = Inches(margin + i * (cw - overlap))
            chev = self._rect(slide, x, y, Inches(cw), h, fill=color,
                              shape=MSO_SHAPE.CHEVRON, shadow=True)
            pad = Inches(0.55) if i > 0 else Inches(0.2)
            tb = self._add_tb(slide, x + pad, y, Inches(cw - 0.7), h)
            blocks = [(label, 15, True, PALETTE["white"])]
            if sub:
                blocks.append((sub, 11, False, PALETTE["cloud"]))
            self._text(tb, blocks, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if caption:
            self._textbox(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.4),
                          [(caption, 14, False, PALETTE["slate"])])
        return slide

    def pipeline_slide(self, header, nodes, kicker=None, caption=None, note=None):
        """Boxes connected by arrows. nodes = list of (title, sub) or (title, sub, color)."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        n = len(nodes)
        margin = 0.55
        arrow_w = 0.5
        total_w = 13.333 - 2 * margin
        bw = (total_w - (n - 1) * arrow_w) / n
        y = Inches(2.75)
        h = Inches(1.7)
        for i, node in enumerate(nodes):
            title = node[0]
            sub = node[1]
            color = node[2] if len(node) > 2 else self.accent
            x = Inches(margin + i * (bw + arrow_w))
            box = self._rect(slide, x, y, Inches(bw), h, fill=PALETTE["white"],
                             line=color, line_w=2.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                             shadow=True)
            box.adjustments[0] = 0.08
            self._rect(slide, x, y, Inches(bw), Inches(0.5), fill=color,
                       shape=TOPCAP)
            self._text(self._add_tb(slide, x, y + Inches(0.04), Inches(bw), Inches(0.42)),
                       [(title, 13, True, PALETTE["white"])],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._text(self._add_tb(slide, x + Inches(0.1), y + Inches(0.62),
                                    Inches(bw - 0.2), Inches(1.0)),
                       [(sub, 11.5, False, PALETTE["slate"])],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < n - 1:
                ax = Inches(margin + (i + 1) * bw + i * arrow_w)
                self._rect(slide, ax, y + Inches(0.6), Inches(arrow_w), Inches(0.5),
                           fill=self.accent, shape=MSO_SHAPE.RIGHT_ARROW)
        if caption:
            self._textbox(slide, Inches(0.6), Inches(4.9), Inches(12.1), Inches(0.7),
                          [(caption, 14, False, PALETTE["slate"])], align=PP_ALIGN.CENTER)
        if note:
            nb = self._rect(slide, Inches(1.4), Inches(5.5), Inches(10.5), Inches(1.0),
                            fill=PALETTE["cloud"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            nb.adjustments[0] = 0.12
            self._text(nb, [(note, 13, False, PALETTE["ink"])],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return slide

    def compare_slide(self, header, left, right, kicker=None):
        """left/right = dict(title=..., lines=[...]). Left is 'problem', right is 'better'."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        panels = [
            (left, Inches(0.6), PALETTE["danger_bg"], PALETTE["danger"]),
            (right, Inches(6.95), PALETTE["green_bg"], PALETTE["green"]),
        ]
        for data, x, bg, edge in panels:
            box = self._rect(slide, x, Inches(1.55), Inches(5.78), Inches(5.25), fill=bg,
                             line=edge, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            box.adjustments[0] = 0.04
            head = self._rect(slide, x, Inches(1.55), Inches(5.78), Inches(0.7), fill=edge,
                              shape=TOPCAP)
            self._text(head, [(data["title"], 16, True, PALETTE["white"])],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            tb = self._add_tb(slide, x + Inches(0.35), Inches(2.5), Inches(5.1), Inches(4.1))
            blocks = []
            for ln in data["lines"]:
                blocks.append(("-  " + ln, 13.5, False, PALETTE["ink"], 10))
            self._text(tb, blocks)
        return slide

    def kpi_slide(self, header, tiles, kicker=None, footer_note=None):
        """tiles = list of (big, label) or (big, label, color)."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        n = len(tiles)
        margin = 0.6
        gap = 0.35
        tw = (13.333 - 2 * margin - (n - 1) * gap) / n
        y = Inches(2.3)
        h = Inches(2.6)
        for i, tile in enumerate(tiles):
            big, label = tile[0], tile[1]
            color = tile[2] if len(tile) > 2 else self.accent
            x = Inches(margin + i * (tw + gap))
            box = self._rect(slide, x, y, Inches(tw), h, fill=PALETTE["white"],
                             line=PALETTE["line"], shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                             shadow=True)
            box.adjustments[0] = 0.06
            self._text(self._add_tb(slide, x, y + Inches(0.35), Inches(tw), Inches(1.2)),
                       [(big, 44, True, color)], align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
            self._text(self._add_tb(slide, x + Inches(0.15), y + Inches(1.6),
                                    Inches(tw - 0.3), Inches(0.9)),
                       [(label, 13, False, PALETTE["slate"])], align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
        if footer_note:
            self._textbox(slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.0),
                          [(footer_note, 14, False, PALETTE["muted"])], align=PP_ALIGN.CENTER)
        return slide

    def chart_slide(self, header, categories, series, kicker=None, caption=None,
                    series_colors=None):
        """Clustered column chart. series = list of (name, [values])."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        cd = CategoryChartData()
        cd.categories = categories
        for name, vals in series:
            cd.add_series(name, vals)
        gx, gy, gw, gh = Inches(1.2), Inches(1.7), Inches(11), Inches(4.4)
        gframe = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, gx, gy, gw, gh, cd)
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(12)
        plot = chart.plots[0]
        plot.gap_width = 80
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size = Pt(12)
        val_ax = chart.value_axis
        val_ax.tick_labels.font.size = Pt(11)
        colors = series_colors or [self.accent, PALETTE["gold"], PALETTE["slate"]]
        for si, plot_series in enumerate(chart.series):
            plot_series.format.fill.solid()
            plot_series.format.fill.fore_color.rgb = colors[si % len(colors)]
        if caption:
            self._textbox(slide, Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.7),
                          [(caption, 13, False, PALETTE["muted"])], align=PP_ALIGN.CENTER)
        return slide

    def timeline_slide(self, header, segments, kicker="2-hour session", labs=None):
        """Horizontal time ribbon. segments = list of (time, label, color?)."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        n = len(segments)
        margin = 0.6
        gap = 0.18
        tw = (13.333 - 2 * margin - (n - 1) * gap) / n
        y = Inches(1.7)
        h = Inches(1.25)
        for i, seg in enumerate(segments):
            time, label = seg[0], seg[1]
            color = seg[2] if len(seg) > 2 else self.accent
            x = Inches(margin + i * (tw + gap))
            box = self._rect(slide, x, y, Inches(tw), h, fill=color,
                             shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
            box.adjustments[0] = 0.09
            self._text(self._add_tb(slide, x, y + Inches(0.12), Inches(tw), Inches(0.4)),
                       [(time, 12, True, PALETTE["gold"])], align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
            self._text(self._add_tb(slide, x + Inches(0.1), y + Inches(0.5),
                                    Inches(tw - 0.2), Inches(0.7)),
                       [(label, 12.5, True, PALETTE["white"])], align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
        if labs:
            self._textbox(slide, Inches(0.6), Inches(3.35), Inches(12), Inches(0.4),
                          [("LABS IN THIS SESSION", 12, True, self.accent)])
            top = 3.8
            for i, lab in enumerate(labs):
                name, desc = (lab if isinstance(lab, tuple) else (lab, ""))
                row = self._rect(slide, Inches(0.6), Inches(top), Inches(12.1), Inches(0.62),
                                 fill=PALETTE["cloud"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                row.adjustments[0] = 0.14
                tag = self._rect(slide, Inches(0.75), Inches(top + 0.11), Inches(1.5),
                                 Inches(0.4), fill=self.accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
                tag.adjustments[0] = 0.3
                self._text(tag, [(name, 11, True, PALETTE["white"])],
                           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                self._text(self._add_tb(slide, Inches(2.45), Inches(top + 0.02),
                                        Inches(10.1), Inches(0.58)),
                           [(desc, 12.5, False, PALETTE["ink"])], anchor=MSO_ANCHOR.MIDDLE)
                top += 0.72
        return slide

    def roadmap_slide(self, header, current_day, days, kicker=None):
        """5-day roadmap; current day highlighted. days = list of (day, label)."""
        slide = self._slide()
        self._chrome(slide, header, kicker)
        n = len(days)
        margin = 0.6
        gap = 0.25
        tw = (13.333 - 2 * margin - (n - 1) * gap) / n
        y = Inches(3.0)
        for i, (day, label) in enumerate(days):
            active = (i + 1) == current_day
            color = self.accent if active else PALETTE["cloud2"]
            txt = PALETTE["white"] if active else PALETTE["muted"]
            h = Inches(1.9) if active else Inches(1.5)
            yy = y - Inches(0.2) if active else y
            x = Inches(margin + i * (tw + gap))
            box = self._rect(slide, x, yy, Inches(tw), h, fill=color,
                             line=PALETTE["line"], shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                             shadow=active)
            box.adjustments[0] = 0.08
            self._text(self._add_tb(slide, x, yy + Inches(0.2), Inches(tw), Inches(0.5)),
                       [(day, 14, True, PALETTE["gold"] if active else self.accent)],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._text(self._add_tb(slide, x + Inches(0.12), yy + Inches(0.7),
                                    Inches(tw - 0.24), Inches(h.inches - 0.8) if hasattr(h, 'inches') else Inches(1.0)),
                       [(label, 11.5, active, txt)], align=PP_ALIGN.CENTER,
                       anchor=MSO_ANCHOR.MIDDLE)
        self._textbox(slide, Inches(0.6), Inches(5.4), Inches(12.1), Inches(0.8),
                      [("You are here: %s" % days[current_day - 1][0], 15, True, self.accent)],
                      align=PP_ALIGN.CENTER)
        return slide

    def takeaways_slide(self, header, points, kicker="Key takeaways"):
        slide = self._slide()
        self._chrome(slide, header, kicker)
        top = 1.7
        for i, pt in enumerate(points):
            row = self._rect(slide, Inches(0.7), Inches(top), Inches(11.9), Inches(0.92),
                             fill=PALETTE["cloud"], line=PALETTE["line"],
                             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            row.adjustments[0] = 0.1
            num = self._rect(slide, Inches(0.95), Inches(top + 0.16), Inches(0.6), Inches(0.6),
                             fill=self.accent, shape=MSO_SHAPE.OVAL)
            self._text(num, [(str(i + 1), 18, True, PALETTE["white"])],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self._text(self._add_tb(slide, Inches(1.8), Inches(top + 0.05),
                                    Inches(10.6), Inches(0.82)),
                       [(pt, 14.5, False, PALETTE["ink"])], anchor=MSO_ANCHOR.MIDDLE)
            top += 1.02
        return slide

    def closing_slide(self, big, sub, next_hint=None):
        slide = self._slide()
        self._rect(slide, 0, 0, SW, SH, fill=PALETTE["ink"])
        self._rect(slide, 0, 0, Inches(0.28), SH, fill=self.accent)
        self._rect(slide, Inches(0.28), 0, Inches(0.08), SH, fill=PALETTE["gold"])
        self._textbox(slide, Inches(1.2), Inches(2.6), Inches(11), Inches(1.2),
                      [(big, 40, True, PALETTE["white"])], align=PP_ALIGN.CENTER)
        self._textbox(slide, Inches(1.6), Inches(3.9), Inches(10.1), Inches(0.9),
                      [(sub, 18, False, PALETTE["line"])], align=PP_ALIGN.CENTER)
        if next_hint:
            chip = self._rect(slide, Inches(4.16), Inches(5.1), Inches(5.0), Inches(0.6),
                              fill=self.accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            chip.adjustments[0] = 0.5
            self._text(chip, [(next_hint, 14, True, PALETTE["white"])],
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return slide

    def save(self, path):
        self.prs.save(path)
        return path
